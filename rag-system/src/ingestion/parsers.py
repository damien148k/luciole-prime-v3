"""
Document Parsers - Multi-format document parsing
Supports: PDF, DOCX, PPTX, XLSX, MSG, EML, TXT
Version améliorée avec ExcelParser pour préserver la structure tabulaire
"""

import os
import signal
import threading
from pathlib import Path
from typing import Dict, List, Optional
from abc import ABC, abstractmethod
from concurrent.futures import ThreadPoolExecutor, TimeoutError as FuturesTimeoutError
from loguru import logger

# Document parsing libraries
import pymupdf
import pymupdf4llm
from docx import Document as DocxDocument
from pptx import Presentation
import openpyxl
import pandas as pd
import extract_msg
from striprtf.striprtf import rtf_to_text
from bs4 import BeautifulSoup

# Import du nouveau ExcelParser amélioré
try:
    from .excel_parser import ExcelParser as AdvancedExcelParser
    ADVANCED_EXCEL_PARSER_AVAILABLE = True
except ImportError:
    ADVANCED_EXCEL_PARSER_AVAILABLE = False
    AdvancedExcelParser = None

# OpenDocument support (optionnel)
try:
    from odf import text as odf_text
    from odf.opendocument import load as odf_load
    from odf.table import Table, TableRow, TableCell
    from odf.draw import Frame
    ODF_AVAILABLE = True
except ImportError:
    ODF_AVAILABLE = False

# Import optionnel pour unstructured (dépendance lourde)
try:
    from unstructured.partition.email import partition_email
    UNSTRUCTURED_AVAILABLE = True
except ImportError:
    UNSTRUCTURED_AVAILABLE = False
    partition_email = None

# Import optionnel pour Apache Tika (fallback .doc et autres formats binaires)
# Nécessite Java installé
try:
    from tika import parser as tika_parser
    TIKA_AVAILABLE = True
except ImportError:
    TIKA_AVAILABLE = False
    tika_parser = None

# Import optionnel pour olefile (extraction .doc sans Java)
try:
    import olefile
    OLEFILE_AVAILABLE = True
except ImportError:
    OLEFILE_AVAILABLE = False
    olefile = None


class BaseParser(ABC):
    """Base class for all document parsers"""
    
    @abstractmethod
    def parse(self, file_path: str) -> Dict:
        """Parse document and return content with metadata"""
        pass
    
    @abstractmethod
    def supported_extensions(self) -> List[str]:
        """Return list of supported file extensions"""
        pass


class PDFParser(BaseParser):
    """
    Parser for PDF documents using PyMuPDF
    Avec support OCR automatique pour les pages scannées ou avec images vectorielles.
    Gère les PDFs vectoriels complexes (plans CAD) avec timeout et fallback.
    """
    
    def __init__(
        self, 
        enable_ocr: bool = True, 
        ocr_languages: List[str] = None,
        min_chars_per_page: int = 50,
        min_words_per_page: int = 8,
        ocr_dpi: int = 300,
        markdown_timeout: int = 300,
        markdown_timeout_per_page: float = 2.0,
        max_drawings_per_page: int = 500
    ):
        """
        Args:
            enable_ocr: Activer l'OCR automatique pour les pages scannées
            ocr_languages: Langues pour l'OCR (défaut: ['fr', 'en'])
            min_chars_per_page: Seuil minimum de caractères par page pour éviter OCR
            min_words_per_page: Seuil minimum de mots par page pour éviter OCR
            ocr_dpi: Résolution DPI pour le rendu des pages (plus haut = meilleure qualité OCR)
            markdown_timeout: Timeout de base en secondes pour pymupdf4llm.to_markdown (défaut: 300s)
            markdown_timeout_per_page: Secondes supplémentaires par page (défaut: 2.0s/page)
            max_drawings_per_page: Seuil de dessins/page au-delà duquel on utilise l'extraction simple
        """
        self.enable_ocr = enable_ocr
        self.ocr_languages = ocr_languages or ['fr', 'en']
        self.min_chars_per_page = min_chars_per_page
        self.min_words_per_page = min_words_per_page
        self.ocr_dpi = ocr_dpi
        self.markdown_timeout = markdown_timeout
        self.markdown_timeout_per_page = markdown_timeout_per_page
        self.max_drawings_per_page = max_drawings_per_page
        self._ocr_processor = None
    
    @property
    def ocr_processor(self):
        """Lazy loading du processeur OCR"""
        if self._ocr_processor is None and self.enable_ocr:
            try:
                from .ocr import OCRProcessor, is_ocr_available
                if is_ocr_available():
                    self._ocr_processor = OCRProcessor(
                        languages=self.ocr_languages,
                        gpu=True
                    )
            except Exception as e:
                logger.warning(f"OCR non disponible: {e}")
        return self._ocr_processor
    
    def supported_extensions(self) -> List[str]:
        return [".pdf"]
    
    def _is_complex_vector_pdf(self, doc) -> bool:
        """
        Détecte si un PDF est un document vectoriel complexe (plan CAD, carte, etc.)
        Ces documents peuvent bloquer pymupdf4llm.to_markdown() pendant très longtemps.
        
        Critères:
        - Beaucoup de dessins/paths vectoriels par page
        - Peu de texte structuré par rapport aux éléments graphiques
        """
        total_drawings = 0
        total_text_blocks = 0
        
        for page in doc:
            # Compter les éléments de dessin (paths vectoriels)
            drawings = page.get_drawings()
            total_drawings += len(drawings)
            
            # Compter les blocs de texte
            text_dict = page.get_text("dict", flags=pymupdf.TEXT_PRESERVE_WHITESPACE)
            for block in text_dict.get("blocks", []):
                if block.get("type") == 0:  # Type 0 = texte
                    total_text_blocks += 1
        
        avg_drawings_per_page = total_drawings / max(1, len(doc))
        
        is_complex = avg_drawings_per_page > self.max_drawings_per_page
        
        if is_complex:
            logger.warning(
                f"PDF vectoriel complexe détecté: {avg_drawings_per_page:.0f} dessins/page "
                f"(seuil: {self.max_drawings_per_page}), {total_text_blocks} blocs texte"
            )
        
        return is_complex
    
    def _extract_text_simple(self, doc) -> str:
        """
        Extraction simple du texte avec PyMuPDF (rapide, sans mise en forme Markdown).
        Utilisé comme fallback pour les PDFs vectoriels complexes.
        """
        text_parts = []
        for page_num, page in enumerate(doc):
            text = page.get_text("text")
            if text.strip():
                text_parts.append(f"--- Page {page_num + 1} ---\n{text}")
        
        return "\n\n".join(text_parts)
    
    def _extract_markdown_pagewise(self, file_path: str, page_count: int, batch_size: int = 10) -> str:
        """
        Extrait le texte Markdown page par page (ou par batch) avec pymupdf4llm.
        Beaucoup plus robuste que l'extraction en une seule fois pour les gros PDF.
        
        Args:
            file_path: Chemin vers le PDF
            page_count: Nombre total de pages
            batch_size: Nombre de pages par batch (défaut: 10)
        """
        all_parts = []
        
        for start in range(0, page_count, batch_size):
            end = min(start + batch_size, page_count)
            pages = list(range(start, end))
            
            try:
                md_text = pymupdf4llm.to_markdown(file_path, pages=pages)
                if md_text and md_text.strip():
                    all_parts.append(md_text)
                logger.debug(f"pymupdf4llm pages {start+1}-{end}/{page_count} OK ({len(md_text)} chars)")
            except Exception as e:
                logger.warning(f"pymupdf4llm échoué pages {start+1}-{end}: {e} - fallback extraction simple")
                # Fallback page par page avec extraction simple pour ce batch
                doc = pymupdf.open(file_path)
                for p in pages:
                    text = doc[p].get_text("text")
                    if text.strip():
                        all_parts.append(f"--- Page {p + 1} ---\n{text}")
                doc.close()
            
            # Log de progression tous les 50 pages
            if end % 50 == 0 or end == page_count:
                logger.info(f"Extraction pymupdf4llm: {end}/{page_count} pages traitées")
        
        return "\n\n".join(all_parts)
    
    def parse(self, file_path: str) -> Dict:
        logger.info(f"Parsing PDF: {file_path}")
        try:
            # Get metadata and check pages first
            doc = pymupdf.open(file_path)
            metadata = doc.metadata
            page_count = len(doc)
            
            # Détecter si c'est un PDF vectoriel complexe (plan CAD, carte, etc.)
            is_complex = self._is_complex_vector_pdf(doc)
            
            # Détecter les pages nécessitant OCR
            pages_need_ocr = []
            ocr_applied = False
            ocr_pages_count = 0
            extraction_method = "pymupdf4llm"
            
            if self.enable_ocr:
                pages_need_ocr = self._detect_scanned_pages(doc)
                if pages_need_ocr:
                    logger.info(f"OCR EasyOCR (GPU) requis pour {len(pages_need_ocr)}/{page_count} pages")
            
            # Choisir la méthode d'extraction
            if is_complex:
                # PDF vectoriel complexe: utiliser extraction simple (rapide)
                logger.info("Utilisation de l'extraction simple pour PDF vectoriel complexe")
                md_text = self._extract_text_simple(doc)
                extraction_method = "simple"
            else:
                # PDF normal: extraction pymupdf4llm page par page (robuste pour gros PDFs)
                doc.close()  # Fermer avant pymupdf4llm
                
                logger.info(f"Extraction pymupdf4llm page par page ({page_count} pages, batch=10)")
                md_text = self._extract_markdown_pagewise(file_path, page_count, batch_size=10)
                
                doc = pymupdf.open(file_path)  # Réouvrir pour OCR si nécessaire
            
            # Vérifier si le texte global est trop court
            total_text_len = len(md_text.strip()) if md_text else 0
            total_words = len(md_text.split()) if md_text else 0
            
            # Si très peu de texte ET OCR disponible, forcer OCR sur toutes les pages
            if self.enable_ocr and self.ocr_processor:
                if total_text_len < self.min_chars_per_page * page_count // 2:
                    logger.warning(f"Texte global insuffisant ({total_text_len} chars, {total_words} mots) - OCR étendu")
                    pages_need_ocr = list(range(page_count))
                
                if pages_need_ocr:
                    ocr_text = self._apply_ocr(file_path, pages_need_ocr)
                    md_text = self._merge_text_with_ocr(md_text or "", ocr_text, pages_need_ocr)
                    ocr_applied = True
                    ocr_pages_count = len([p for p in pages_need_ocr if ocr_text.get(p, "").strip()])
                    logger.info(f"OCR terminé: {ocr_pages_count} pages avec texte extrait")
            
            if hasattr(doc, 'close') and not doc.is_closed:
                doc.close()
            
            logger.info(f"PDF parsé: {page_count} pages, {total_text_len} chars, méthode={extraction_method}")
            
            return {
                "content": md_text or "",
                "metadata": {
                    "type": "pdf",
                    "file_name": Path(file_path).name,
                    "file_path": file_path,
                    "page_count": page_count,
                    "title": metadata.get("title", ""),
                    "author": metadata.get("author", ""),
                    "creation_date": metadata.get("creationDate", ""),
                    "ocr_applied": ocr_applied,
                    "ocr_pages_count": ocr_pages_count,
                    "extraction_method": extraction_method,
                    "is_complex_vector": is_complex
                }
            }
        except Exception as e:
            logger.error(f"Error parsing PDF {file_path}: {e}")
            raise
    
    def _page_needs_ocr(self, text: str) -> bool:
        """
        Détermine si une page nécessite l'OCR basé sur des seuils de caractères ET de mots.
        Cette double vérification est importante pour les PDF avec images vectorielles
        qui peuvent avoir peu de texte extractible mais des éléments graphiques.
        """
        text = text.strip()
        char_count = len(text)
        word_count = len(text.split())
        return char_count < self.min_chars_per_page or word_count < self.min_words_per_page
    
    def _detect_scanned_pages(self, doc) -> List[int]:
        """
        Détecte les pages qui semblent être scannées ou contenir principalement
        des images vectorielles (peu de texte extractible).
        
        Utilise un double seuil: caractères ET mots pour une détection plus robuste.
        """
        pages_need_ocr = []
        for page_num, page in enumerate(doc):
            text = page.get_text().strip()
            if self._page_needs_ocr(text):
                pages_need_ocr.append(page_num)
                logger.debug(f"Page {page_num + 1}: OCR requis (chars={len(text)}, words={len(text.split())})")
        return pages_need_ocr
    
    def _apply_ocr(self, file_path: str, pages: List[int]) -> Dict[int, str]:
        """
        Applique l'OCR EasyOCR (GPU) sur les pages spécifiées.
        Utilise le DPI configuré pour un meilleur rendu des images vectorielles.
        """
        ocr_results = {}
        for page_num in pages:
            try:
                text = self.ocr_processor.ocr_pdf_page(file_path, page_num, dpi=self.ocr_dpi)
                ocr_results[page_num] = text
                logger.debug(f"OCR page {page_num + 1}: {len(text)} caractères extraits (DPI={self.ocr_dpi})")
            except Exception as e:
                logger.warning(f"OCR échoué page {page_num + 1}: {e}")
                ocr_results[page_num] = ""
        return ocr_results
    
    def _merge_text_with_ocr(self, original_text: str, ocr_text: Dict[int, str], pages: List[int]) -> str:
        """Fusionne le texte original avec le texte OCR"""
        if not ocr_text:
            return original_text
        
        # Ajouter le texte OCR à la fin avec indication des pages
        ocr_sections = []
        for page_num in sorted(ocr_text.keys()):
            if ocr_text[page_num].strip():
                ocr_sections.append(f"\n--- OCR Page {page_num + 1} ---\n{ocr_text[page_num]}")
        
        if ocr_sections:
            return original_text + "\n\n" + "\n".join(ocr_sections)
        
        return original_text


class DOCXParser(BaseParser):
    """
    Parser for Microsoft Word documents (.docx et .doc)
    
    - .docx : utilise python-docx (natif)
    - .doc : utilise olefile (extraction texte du format OLE) ou fallback Tika
    """
    
    def supported_extensions(self) -> List[str]:
        return [".docx", ".doc"]
    
    def _parse_doc_with_olefile(self, file_path: str) -> Dict:
        """
        Parse un fichier .doc (format OLE/binaire) avec olefile.
        Extraction du texte depuis le stream 'WordDocument'.
        Ne nécessite pas Java.
        """
        if not OLEFILE_AVAILABLE:
            raise RuntimeError("olefile non disponible pour parser les fichiers .doc")
        
        logger.info(f"Parsing DOC avec olefile: {file_path}")
        
        try:
            ole = olefile.OleFileIO(file_path)
            
            # Chercher le stream de texte Word
            text_parts = []
            
            # Méthode 1: Extraire depuis "WordDocument" stream (texte brut encodé)
            if ole.exists('WordDocument'):
                word_stream = ole.openstream('WordDocument').read()
                # Le texte est souvent après un header, on cherche les caractères ASCII
                text = word_stream.decode('latin-1', errors='ignore')
                # Filtrer pour garder uniquement le texte lisible
                clean_text = ''.join(
                    c if c.isprintable() or c in '\n\r\t' else ' ' 
                    for c in text
                )
                text_parts.append(clean_text)
            
            # Méthode 2: Chercher dans d'autres streams courants
            for stream_name in ['1Table', '0Table', 'Data']:
                if ole.exists(stream_name):
                    try:
                        stream_data = ole.openstream(stream_name).read()
                        decoded = stream_data.decode('utf-16-le', errors='ignore')
                        if decoded.strip():
                            text_parts.append(decoded)
                    except Exception:
                        pass
            
            ole.close()
            
            # Combiner et nettoyer
            content = '\n'.join(text_parts)
            # Supprimer les lignes avec trop peu de caractères (bruit)
            lines = [line.strip() for line in content.split('\n') 
                     if line.strip() and len(line.strip()) > 2]
            content = '\n'.join(lines)
            
            # Nettoyer les espaces multiples
            import re
            content = re.sub(r' +', ' ', content)
            content = re.sub(r'\n+', '\n', content)
            
            if not content.strip():
                raise ValueError("Aucun texte extractible avec olefile")
            
            return {
                "content": content,
                "metadata": {
                    "type": "doc",
                    "file_name": Path(file_path).name,
                    "file_path": file_path,
                    "parser": "olefile",
                    "char_count": len(content)
                }
            }
        except Exception as e:
            logger.warning(f"olefile parsing failed for {file_path}: {e}")
            raise
    
    def _parse_with_tika(self, file_path: str) -> Dict:
        """
        Parse un document avec Apache Tika (fallback pour .doc binaires)
        Nécessite Java installé sur le système.
        """
        if not TIKA_AVAILABLE:
            raise RuntimeError(
                "Apache Tika non disponible. Installez-le avec: pip install tika\n"
                "Note: Java doit être installé sur le système."
            )
        
        logger.info(f"Parsing avec Tika (fallback .doc): {file_path}")
        
        try:
            parsed = tika_parser.from_file(file_path)
            content = parsed.get('content', '') or ''
            metadata_tika = parsed.get('metadata', {}) or {}
            
            # Nettoyer le contenu (Tika peut ajouter beaucoup d'espaces)
            content = '\n'.join(line.strip() for line in content.split('\n') if line.strip())
            
            return {
                "content": content,
                "metadata": {
                    "type": "doc",
                    "file_name": Path(file_path).name,
                    "file_path": file_path,
                    "parser": "tika",
                    "title": metadata_tika.get('title', ''),
                    "author": metadata_tika.get('Author', '') or metadata_tika.get('creator', ''),
                    "char_count": len(content)
                }
            }
        except Exception as e:
            logger.error(f"Tika parsing failed for {file_path}: {e}")
            raise
    
    def parse(self, file_path: str) -> Dict:
        ext = Path(file_path).suffix.lower()
        
        # Pour les .doc binaires, essayer olefile d'abord, puis Tika
        if ext == ".doc":
            logger.info(f"Parsing DOC (ancien format): {file_path}")
            
            # Essayer olefile (pas besoin de Java)
            if OLEFILE_AVAILABLE:
                try:
                    return self._parse_doc_with_olefile(file_path)
                except Exception as e:
                    logger.warning(f"olefile failed: {e}")
            
            # Fallback vers Tika
            if TIKA_AVAILABLE:
                return self._parse_with_tika(file_path)
            
            raise RuntimeError(
                f"Impossible de parser le fichier .doc: {file_path}\n"
                "Installez olefile (pip install olefile) ou tika (pip install tika + Java)"
            )
        
        # Pour les .docx, essayer python-docx d'abord
        logger.info(f"Parsing DOCX: {file_path}")
        try:
            doc = DocxDocument(file_path)
            
            # Extract text from paragraphs
            paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
            content = "\n\n".join(paragraphs)
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text for cell in row.cells)
                    content += f"\n{row_text}"
            
            return {
                "content": content,
                "metadata": {
                    "type": "docx",
                    "file_name": Path(file_path).name,
                    "file_path": file_path,
                    "paragraph_count": len(paragraphs),
                    "table_count": len(doc.tables),
                    "parser": "python-docx"
                }
            }
        except Exception as e:
            # Si python-docx échoue (fichier corrompu ou mal nommé), essayer les fallbacks
            logger.warning(f"python-docx failed for {file_path}: {e}")
            
            # Essayer olefile
            if OLEFILE_AVAILABLE:
                try:
                    logger.info("Fallback vers olefile...")
                    return self._parse_doc_with_olefile(file_path)
                except Exception:
                    pass
            
            # Essayer Tika
            if TIKA_AVAILABLE:
                logger.info("Fallback vers Tika...")
                return self._parse_with_tika(file_path)
            
            logger.error(f"Error parsing DOCX {file_path}: {e}")
            raise


class PPTXParser(BaseParser):
    """
    Parser for PowerPoint presentations (.pptx et .ppt)
    
    - .pptx : utilise python-pptx (natif)
    - .ppt : utilise olefile (extraction texte du format OLE) ou fallback Tika
    """
    
    def supported_extensions(self) -> List[str]:
        return [".pptx", ".ppt"]
    
    def _parse_ppt_with_olefile(self, file_path: str) -> Dict:
        """
        Parse un fichier .ppt (format OLE/binaire) avec olefile.
        Extraction du texte depuis les streams PowerPoint.
        Ne nécessite pas Java.
        """
        if not OLEFILE_AVAILABLE:
            raise RuntimeError("olefile non disponible pour parser les fichiers .ppt")
        
        logger.info(f"Parsing PPT avec olefile: {file_path}")
        
        try:
            ole = olefile.OleFileIO(file_path)
            
            text_parts = []
            
            # PowerPoint stocke le texte dans plusieurs streams possibles
            ppt_streams = [
                'PowerPoint Document',
                'Current User',
                'Pictures'
            ]
            
            for stream_name in ppt_streams:
                if ole.exists(stream_name):
                    try:
                        stream_data = ole.openstream(stream_name).read()
                        # Essayer plusieurs encodages
                        for encoding in ['utf-16-le', 'latin-1', 'utf-8']:
                            try:
                                decoded = stream_data.decode(encoding, errors='ignore')
                                # Filtrer les caractères lisibles
                                clean = ''.join(
                                    c if c.isprintable() or c in '\n\r\t' else ' '
                                    for c in decoded
                                )
                                if clean.strip():
                                    text_parts.append(clean)
                                break
                            except Exception:
                                continue
                    except Exception:
                        pass
            
            ole.close()
            
            # Combiner et nettoyer
            content = '\n'.join(text_parts)
            import re
            content = re.sub(r' +', ' ', content)
            content = re.sub(r'\n+', '\n', content)
            
            # Filtrer les lignes trop courtes (bruit)
            lines = [line.strip() for line in content.split('\n') 
                     if line.strip() and len(line.strip()) > 3]
            content = '\n'.join(lines)
            
            if not content.strip():
                raise ValueError("Aucun texte extractible avec olefile")
            
            return {
                "content": content,
                "metadata": {
                    "type": "ppt",
                    "file_name": Path(file_path).name,
                    "file_path": file_path,
                    "parser": "olefile",
                    "char_count": len(content)
                }
            }
        except Exception as e:
            logger.warning(f"olefile parsing failed for {file_path}: {e}")
            raise
    
    def _parse_with_tika(self, file_path: str) -> Dict:
        """Parse un fichier .ppt avec Apache Tika"""
        if not TIKA_AVAILABLE:
            raise RuntimeError(
                "Apache Tika non disponible pour les fichiers .ppt\n"
                "Installez-le avec: pip install tika (Java requis)"
            )
        
        logger.info(f"Parsing PPT avec Tika: {file_path}")
        
        try:
            parsed = tika_parser.from_file(file_path)
            content = parsed.get('content', '') or ''
            
            # Nettoyer le contenu
            content = '\n'.join(line.strip() for line in content.split('\n') if line.strip())
            
            # Essayer de détecter les slides via les séparateurs Tika
            # Tika insère souvent des lignes vides entre les slides
            slides = [s.strip() for s in content.split('\n\n\n') if s.strip()]
            
            return {
                "content": content,
                "metadata": {
                    "type": "ppt",
                    "file_name": Path(file_path).name,
                    "file_path": file_path,
                    "slide_count": len(slides) if slides else 1,
                    "parser": "tika"
                }
            }
        except Exception as e:
            logger.error(f"Tika parsing failed for {file_path}: {e}")
            raise
    
    def parse(self, file_path: str) -> Dict:
        ext = Path(file_path).suffix.lower()
        
        # Pour les .ppt binaires, essayer olefile puis Tika
        if ext == ".ppt":
            logger.info(f"Parsing PPT (ancien format): {file_path}")
            
            # Essayer olefile d'abord (pas besoin de Java)
            if OLEFILE_AVAILABLE:
                try:
                    return self._parse_ppt_with_olefile(file_path)
                except Exception as e:
                    logger.warning(f"olefile failed: {e}")
            
            # Fallback vers Tika
            if TIKA_AVAILABLE:
                return self._parse_with_tika(file_path)
            
            raise RuntimeError(
                f"Impossible de parser le fichier .ppt: {file_path}\n"
                "Installez olefile (pip install olefile) ou tika (pip install tika + Java)"
            )
        
        # Pour les .pptx, utiliser python-pptx
        logger.info(f"Parsing PPTX: {file_path}")
        try:
            prs = Presentation(file_path)
            slides_content = []
            
            for i, slide in enumerate(prs.slides, 1):
                slide_text = f"--- Slide {i} ---\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text += shape.text + "\n"
                slides_content.append(slide_text)
            
            return {
                "content": "\n\n".join(slides_content),
                "metadata": {
                    "type": "pptx",
                    "file_name": Path(file_path).name,
                    "file_path": file_path,
                    "slide_count": len(prs.slides),
                    "parser": "python-pptx"
                }
            }
        except Exception as e:
            # Fallback vers olefile puis Tika
            logger.warning(f"python-pptx failed for {file_path}: {e}")
            
            if OLEFILE_AVAILABLE:
                try:
                    logger.info("Fallback vers olefile...")
                    return self._parse_ppt_with_olefile(file_path)
                except Exception:
                    pass
            
            if TIKA_AVAILABLE:
                logger.info("Fallback vers Tika...")
                return self._parse_with_tika(file_path)
            
            logger.error(f"Error parsing PPTX {file_path}: {e}")
            raise


class XLSXParser(BaseParser):
    """
    Parser for Excel spreadsheets
    Version améliorée: génère du Markdown structuré + stockage SQLite optionnel
    
    Si AdvancedExcelParser est disponible, l'utilise pour:
    - Préserver la structure tabulaire (Markdown)
    - Stocker les données brutes en SQLite
    - Calculer des statistiques
    
    Sinon, fallback vers parsing basique.
    """
    
    def __init__(
        self,
        max_rows_per_chunk: int = 50,
        enable_sql_storage: bool = True,
        sqlite_path: str = "/app/data/excel_data.db"
    ):
        """
        Initialise le parser Excel
        
        Args:
            max_rows_per_chunk: Lignes max par chunk Markdown
            enable_sql_storage: Activer stockage SQLite
            sqlite_path: Chemin vers la base SQLite
        """
        self.max_rows_per_chunk = max_rows_per_chunk
        self.enable_sql_storage = enable_sql_storage
        self.sqlite_path = sqlite_path
        
        # Utiliser le parser avancé si disponible
        self._advanced_parser = None
        if ADVANCED_EXCEL_PARSER_AVAILABLE and AdvancedExcelParser:
            try:
                self._advanced_parser = AdvancedExcelParser(
                    max_rows_per_chunk=max_rows_per_chunk,
                    enable_sql_storage=enable_sql_storage,
                    sqlite_path=sqlite_path
                )
                logger.info("ExcelParser avancé initialisé (Markdown + SQLite)")
            except Exception as e:
                logger.warning(f"Impossible d'initialiser ExcelParser avancé: {e}")
    
    def supported_extensions(self) -> List[str]:
        return [".xlsx", ".xls"]
    
    def parse(self, file_path: str) -> Dict:
        logger.info(f"Parsing XLSX: {file_path}")
        
        # Utiliser le parser avancé si disponible
        if self._advanced_parser:
            try:
                return self._advanced_parser.parse(file_path)
            except Exception as e:
                logger.warning(f"ExcelParser avancé échoué, fallback: {e}")
        
        # Fallback: parsing basique avec Markdown
        return self._parse_basic_markdown(file_path)
    
    def _parse_basic_markdown(self, file_path: str) -> Dict:
        """Parsing basique avec conversion Markdown (fallback)"""
        try:
            excel_file = pd.ExcelFile(file_path)
            sheets_content = []
            sheet_metadata = []
            
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(excel_file, sheet_name=sheet_name)
                
                if df.empty:
                    continue
                
                # Générer Markdown au lieu de to_string()
                sheet_text = f"## Tableau: {sheet_name}\n\n"
                sheet_text += f"**Colonnes**: {', '.join(str(c) for c in df.columns)}\n"
                sheet_text += f"**Lignes**: {len(df)}\n\n"
                
                # Table Markdown
                try:
                    sheet_text += df.head(self.max_rows_per_chunk).to_markdown(index=False)
                except Exception:
                    # Fallback manuel
                    sheet_text += self._manual_markdown(df.head(self.max_rows_per_chunk))
                
                if len(df) > self.max_rows_per_chunk:
                    sheet_text += f"\n\n*... ({len(df) - self.max_rows_per_chunk} lignes supplémentaires)*"
                
                sheets_content.append(sheet_text)
                sheet_metadata.append({
                    "sheet_name": sheet_name,
                    "columns": list(df.columns),
                    "row_count": len(df)
                })
            
            excel_file.close()
            
            return {
                "content": "\n\n---\n\n".join(sheets_content),
                "metadata": {
                    "type": "excel",
                    "file_name": Path(file_path).name,
                    "file_path": file_path,
                    "sheet_count": len(excel_file.sheet_names),
                    "sheets": sheet_metadata
                }
            }
        except Exception as e:
            logger.error(f"Error parsing XLSX {file_path}: {e}")
            raise
    
    def _manual_markdown(self, df: pd.DataFrame) -> str:
        """Génère un tableau Markdown manuellement"""
        cols = list(df.columns)
        header = "| " + " | ".join(str(c) for c in cols) + " |"
        separator = "|" + "|".join(["---"] * len(cols)) + "|"
        
        rows = []
        for _, row in df.iterrows():
            values = [str(v) if pd.notna(v) else "" for v in row.values]
            rows.append("| " + " | ".join(values) + " |")
        
        return header + "\n" + separator + "\n" + "\n".join(rows)


class MSGParser(BaseParser):
    """Parser for Outlook MSG files"""
    
    def supported_extensions(self) -> List[str]:
        return [".msg"]
    
    def parse(self, file_path: str) -> Dict:
        logger.info(f"Parsing MSG: {file_path}")
        try:
            msg = extract_msg.Message(file_path)
            
            content = f"Subject: {msg.subject}\n"
            content += f"From: {msg.sender}\n"
            content += f"To: {msg.to}\n"
            content += f"Date: {msg.date}\n\n"
            content += msg.body or ""
            
            return {
                "content": content,
                "metadata": {
                    "type": "msg",
                    "file_name": Path(file_path).name,
                    "file_path": file_path,
                    "subject": msg.subject,
                    "sender": msg.sender,
                    "to": msg.to,
                    "date": str(msg.date),
                    "has_attachments": len(msg.attachments) > 0
                }
            }
        except Exception as e:
            logger.error(f"Error parsing MSG {file_path}: {e}")
            raise


class EMLParser(BaseParser):
    """Parser for EML email files"""
    
    def supported_extensions(self) -> List[str]:
        return [".eml"]
    
    def parse(self, file_path: str) -> Dict:
        logger.info(f"Parsing EML: {file_path}")
        
        if not UNSTRUCTURED_AVAILABLE:
            # Fallback: lecture basique du fichier EML
            logger.warning("unstructured not available, using basic EML parsing")
            import email
            from email import policy
            
            with open(file_path, "rb") as f:
                msg = email.message_from_binary_file(f, policy=policy.default)
            
            # Extraire le contenu
            body = ""
            if msg.is_multipart():
                for part in msg.walk():
                    if part.get_content_type() == "text/plain":
                        body += part.get_content()
            else:
                body = msg.get_content()
            
            return {
                "content": f"Subject: {msg['subject']}\nFrom: {msg['from']}\nTo: {msg['to']}\n\n{body}",
                "metadata": {
                    "type": "email",
                    "file_name": Path(file_path).name,
                    "file_path": file_path,
                    "subject": msg["subject"],
                    "sender": msg["from"],
                    "to": msg["to"]
                }
            }
        
        try:
            elements = partition_email(filename=file_path)
            
            content = "\n".join([el.text for el in elements if el.text])
            
            # Extract metadata from elements
            metadata = {
                "type": "email",
                "file_name": Path(file_path).name,
                "file_path": file_path
            }
            
            for el in elements:
                if hasattr(el, "metadata"):
                    if hasattr(el.metadata, "subject"):
                        metadata["subject"] = el.metadata.subject
                    if hasattr(el.metadata, "sent_from"):
                        metadata["sender"] = el.metadata.sent_from
                    if hasattr(el.metadata, "sent_to"):
                        metadata["to"] = el.metadata.sent_to
            
            return {
                "content": content,
                "metadata": metadata
            }
        except Exception as e:
            logger.error(f"Error parsing EML {file_path}: {e}")
            raise


class ImageParser(BaseParser):
    """Parser for images using OCR (PNG, JPG, JPEG, TIFF, BMP)"""
    
    def __init__(self, ocr_languages: List[str] = None):
        self.ocr_languages = ocr_languages or ['fr', 'en']
        self._ocr_processor = None
    
    @property
    def ocr_processor(self):
        """Lazy loading du processeur OCR"""
        if self._ocr_processor is None:
            try:
                from .ocr import OCRProcessor, is_ocr_available
                if is_ocr_available():
                    self._ocr_processor = OCRProcessor(
                        languages=self.ocr_languages,
                        gpu=True
                    )
            except Exception as e:
                logger.warning(f"OCR non disponible: {e}")
        return self._ocr_processor
    
    def supported_extensions(self) -> List[str]:
        return [".png", ".jpg", ".jpeg", ".tiff", ".tif", ".bmp"]
    
    def parse(self, file_path: str) -> Dict:
        logger.info(f"Parsing Image with OCR: {file_path}")
        
        if not self.ocr_processor:
            raise RuntimeError("OCR non disponible. Installez EasyOCR: pip install easyocr")
        
        try:
            content = self.ocr_processor.ocr_image(file_path)
            
            return {
                "content": content,
                "metadata": {
                    "type": "image",
                    "file_name": Path(file_path).name,
                    "file_path": file_path,
                    "ocr_applied": True
                }
            }
        except Exception as e:
            logger.error(f"Error parsing image {file_path}: {e}")
            raise


class TXTParser(BaseParser):
    """Parser for plain text files"""
    
    def supported_extensions(self) -> List[str]:
        return [".txt", ".md", ".rst", ".csv", ".log", ".ini", ".cfg", ".json", ".yaml", ".yml"]
    
    def parse(self, file_path: str) -> Dict:
        logger.info(f"Parsing TXT: {file_path}")
        try:
            # Essayer plusieurs encodages
            content = None
            for encoding in ["utf-8", "latin-1", "cp1252", "iso-8859-1"]:
                try:
                    with open(file_path, "r", encoding=encoding) as f:
                        content = f.read()
                    break
                except UnicodeDecodeError:
                    continue
            
            if content is None:
                # Fallback: lecture binaire
                with open(file_path, "rb") as f:
                    content = f.read().decode("utf-8", errors="ignore")
            
            return {
                "content": content,
                "metadata": {
                    "type": "text",
                    "file_name": Path(file_path).name,
                    "file_path": file_path,
                    "char_count": len(content),
                    "line_count": content.count("\n") + 1
                }
            }
        except Exception as e:
            logger.error(f"Error parsing TXT {file_path}: {e}")
            raise


class RTFParser(BaseParser):
    """Parser for Rich Text Format files"""
    
    def supported_extensions(self) -> List[str]:
        return [".rtf"]
    
    def parse(self, file_path: str) -> Dict:
        logger.info(f"Parsing RTF: {file_path}")
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                rtf_content = f.read()
            
            # Convertir RTF en texte brut
            content = rtf_to_text(rtf_content)
            
            return {
                "content": content,
                "metadata": {
                    "type": "rtf",
                    "file_name": Path(file_path).name,
                    "file_path": file_path,
                    "char_count": len(content)
                }
            }
        except Exception as e:
            logger.error(f"Error parsing RTF {file_path}: {e}")
            raise


class HTMLParser(BaseParser):
    """Parser for HTML/HTM web pages"""
    
    def supported_extensions(self) -> List[str]:
        return [".html", ".htm", ".xhtml"]
    
    def parse(self, file_path: str) -> Dict:
        logger.info(f"Parsing HTML: {file_path}")
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                html_content = f.read()
            
            # Parser avec BeautifulSoup
            soup = BeautifulSoup(html_content, "html.parser")
            
            # Supprimer les scripts et styles
            for script in soup(["script", "style", "meta", "link"]):
                script.decompose()
            
            # Extraire le texte
            content = soup.get_text(separator="\n", strip=True)
            
            # Extraire le titre si present
            title = ""
            if soup.title:
                title = soup.title.string or ""
            
            return {
                "content": content,
                "metadata": {
                    "type": "html",
                    "file_name": Path(file_path).name,
                    "file_path": file_path,
                    "title": title,
                    "char_count": len(content)
                }
            }
        except Exception as e:
            logger.error(f"Error parsing HTML {file_path}: {e}")
            raise


class ODFParser(BaseParser):
    """Parser for OpenDocument Format files (ODT, ODS, ODP)"""
    
    def supported_extensions(self) -> List[str]:
        return [".odt", ".ods", ".odp"]
    
    def parse(self, file_path: str) -> Dict:
        ext = Path(file_path).suffix.lower()
        logger.info(f"Parsing ODF ({ext}): {file_path}")
        
        if not ODF_AVAILABLE:
            raise RuntimeError("odfpy non installe. Executez: pip install odfpy")
        
        try:
            doc = odf_load(file_path)
            
            if ext == ".odt":
                return self._parse_odt(doc, file_path)
            elif ext == ".ods":
                return self._parse_ods(doc, file_path)
            elif ext == ".odp":
                return self._parse_odp(doc, file_path)
            else:
                raise ValueError(f"Extension ODF non supportee: {ext}")
                
        except Exception as e:
            logger.error(f"Error parsing ODF {file_path}: {e}")
            raise
    
    def _parse_odt(self, doc, file_path: str) -> Dict:
        """Parser pour documents texte OpenDocument"""
        content_parts = []
        
        # Extraire tout le texte des paragraphes
        for element in doc.getElementsByType(odf_text.P):
            text = self._get_text_content(element)
            if text.strip():
                content_parts.append(text)
        
        # Extraire les titres
        for element in doc.getElementsByType(odf_text.H):
            text = self._get_text_content(element)
            if text.strip():
                content_parts.append(f"# {text}")
        
        content = "\n\n".join(content_parts)
        
        return {
            "content": content,
            "metadata": {
                "type": "odt",
                "file_name": Path(file_path).name,
                "file_path": file_path,
                "paragraph_count": len(content_parts)
            }
        }
    
    def _parse_ods(self, doc, file_path: str) -> Dict:
        """Parser pour feuilles de calcul OpenDocument"""
        content_parts = []
        
        for table in doc.getElementsByType(Table):
            table_name = table.getAttribute("name") or "Sheet"
            content_parts.append(f"--- {table_name} ---")
            
            for row in table.getElementsByType(TableRow):
                cells = []
                for cell in row.getElementsByType(TableCell):
                    cell_text = self._get_text_content(cell)
                    cells.append(cell_text)
                if any(cells):
                    content_parts.append(" | ".join(cells))
        
        content = "\n".join(content_parts)
        
        return {
            "content": content,
            "metadata": {
                "type": "ods",
                "file_name": Path(file_path).name,
                "file_path": file_path,
                "table_count": len(doc.getElementsByType(Table))
            }
        }
    
    def _parse_odp(self, doc, file_path: str) -> Dict:
        """Parser pour presentations OpenDocument"""
        content_parts = []
        slide_num = 0
        
        # Les slides sont dans des frames
        for frame in doc.getElementsByType(Frame):
            slide_num += 1
            slide_content = [f"--- Slide {slide_num} ---"]
            
            for p in frame.getElementsByType(odf_text.P):
                text = self._get_text_content(p)
                if text.strip():
                    slide_content.append(text)
            
            if len(slide_content) > 1:
                content_parts.append("\n".join(slide_content))
        
        # Fallback: extraire tous les paragraphes
        if not content_parts:
            for element in doc.getElementsByType(odf_text.P):
                text = self._get_text_content(element)
                if text.strip():
                    content_parts.append(text)
        
        content = "\n\n".join(content_parts)
        
        return {
            "content": content,
            "metadata": {
                "type": "odp",
                "file_name": Path(file_path).name,
                "file_path": file_path,
                "slide_count": slide_num if slide_num > 0 else len(content_parts)
            }
        }
    
    def _get_text_content(self, element) -> str:
        """Extrait le contenu texte d'un element ODF"""
        text_parts = []
        
        if hasattr(element, "childNodes"):
            for child in element.childNodes:
                if hasattr(child, "data"):
                    text_parts.append(child.data)
                elif hasattr(child, "childNodes"):
                    text_parts.append(self._get_text_content(child))
        
        return "".join(text_parts)


class DocumentParser:
    """
    Factory class for document parsing
    Automatically selects the appropriate parser based on file extension
    """
    
    def __init__(self, pdf_config: Dict = None):
        """
        Args:
            pdf_config: Configuration PDF optionnelle (depuis settings.yaml section 'pdf')
                - markdown_timeout: Timeout de base en secondes (défaut: 300)
                - markdown_timeout_per_page: Secondes par page (défaut: 2.0)
                - ocr_dpi: Résolution DPI pour OCR (défaut: 300)
                - max_drawings_per_page: Seuil dessins/page (défaut: 500)
        """
        self.pdf_config = pdf_config or {}
        self.parsers: Dict[str, BaseParser] = {}
        self._register_parsers()
    
    def _register_parsers(self):
        """Register all available parsers"""
        # Extraire la config PDF (avec valeurs par défaut)
        pdf_timeout = self.pdf_config.get("markdown_timeout", 300)
        pdf_timeout_per_page = self.pdf_config.get("markdown_timeout_per_page", 2.0)
        pdf_ocr_dpi = self.pdf_config.get("ocr_dpi", 300)
        pdf_max_drawings = self.pdf_config.get("max_drawings_per_page", 500)
        pdf_enable_ocr = self.pdf_config.get("enable_ocr", True)
        
        parser_classes = [
            # PDF avec OCR EasyOCR (GPU) automatique pour pages scannées/vectorielles
            # Gère les PDFs vectoriels complexes (plans CAD) avec timeout et fallback
            PDFParser(
                enable_ocr=pdf_enable_ocr,
                ocr_languages=['fr', 'en'],
                min_chars_per_page=50,      # Seuil caractères (comme ancien ingest.py)
                min_words_per_page=8,       # Seuil mots (comme ancien ingest.py)
                ocr_dpi=pdf_ocr_dpi,
                markdown_timeout=pdf_timeout,
                markdown_timeout_per_page=pdf_timeout_per_page,
                max_drawings_per_page=pdf_max_drawings
            ),
            DOCXParser(),
            PPTXParser(),
            XLSXParser(),
            MSGParser(),
            EMLParser(),
            TXTParser(),
            ImageParser(),  # Images avec OCR
            RTFParser(),    # Rich Text Format
            HTMLParser(),   # HTML/HTM web pages
            ODFParser()     # OpenDocument (ODT, ODS, ODP)
        ]
        
        for parser in parser_classes:
            for ext in parser.supported_extensions():
                self.parsers[ext.lower()] = parser
    
    def get_supported_extensions(self) -> List[str]:
        """Return list of all supported file extensions"""
        return list(self.parsers.keys())
    
    def parse(self, file_path: str) -> Dict:
        """
        Parse a document and return its content with metadata
        
        Args:
            file_path: Path to the document
            
        Returns:
            Dict with 'content' and 'metadata' keys
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")
        
        ext = Path(file_path).suffix.lower()
        
        if ext not in self.parsers:
            raise ValueError(f"Unsupported file format: {ext}. Supported: {self.get_supported_extensions()}")
        
        parser = self.parsers[ext]
        return parser.parse(file_path)
    
    def parse_directory(self, dir_path: str, recursive: bool = True) -> List[Dict]:
        """
        Parse all supported documents in a directory
        
        Args:
            dir_path: Path to directory
            recursive: Whether to search subdirectories
            
        Returns:
            List of parsed documents
        """
        results = []
        path = Path(dir_path)
        
        pattern = "**/*" if recursive else "*"
        
        for file_path in path.glob(pattern):
            if file_path.is_file() and file_path.suffix.lower() in self.parsers:
                try:
                    result = self.parse(str(file_path))
                    results.append(result)
                except Exception as e:
                    logger.warning(f"Failed to parse {file_path}: {e}")
        
        return results

